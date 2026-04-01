#!/usr/bin/env python3
"""
Sales Comparison Presentation Generator
February vs March 2026 - SUNA BIM Sales Analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os

class SalesPresentation:
    """Professional sales comparison presentation builder"""

    # Color Palette - Professional Light Theme
    COLORS = {
        'primary_blue': RGBColor(41, 128, 185),      # Professional blue
        'secondary_blue': RGBColor(52, 152, 219),    # Lighter blue
        'accent_green': RGBColor(39, 174, 96),       # Success green
        'accent_red': RGBColor(231, 76, 60),         # Alert red
        'dark_gray': RGBColor(52, 73, 94),           # Text
        'light_gray': RGBColor(236, 240, 241),       # Background
        'white': RGBColor(255, 255, 255),
        'medium_gray': RGBColor(149, 165, 166)       # Secondary text
    }

    def __init__(self):
        """Initialize presentation with 16:9 layout"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)

    def add_title_slide(self):
        """Create professional title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.prs.slide_width,
            self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS['white']
        background.line.fill.background()

        # Top accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.prs.slide_width,
            Inches(0.3)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.COLORS['primary_blue']
        accent_bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5),
            Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Sales Performance Analysis"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['primary_blue']
        title_para.alignment = PP_ALIGN.CENTER

        # Thai subtitle
        thai_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.3),
            Inches(8), Inches(0.5)
        )
        thai_frame = thai_box.text_frame
        thai_frame.text = "(การวิเคราะห์ผลการดำเนินงานด้านยอดขาย)"
        thai_para = thai_frame.paragraphs[0]
        thai_para.font.size = Pt(20)
        thai_para.font.color.rgb = self.COLORS['medium_gray']
        thai_para.alignment = PP_ALIGN.CENTER

        # Period
        period_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.2),
            Inches(8), Inches(0.6)
        )
        period_frame = period_box.text_frame
        period_frame.text = "February vs March 2026"
        period_para = period_frame.paragraphs[0]
        period_para.font.size = Pt(28)
        period_para.font.color.rgb = self.COLORS['dark_gray']
        period_para.alignment = PP_ALIGN.CENTER

        # Thai period
        period_thai_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.7),
            Inches(8), Inches(0.4)
        )
        period_thai_frame = period_thai_box.text_frame
        period_thai_frame.text = "(กุมภาพันธ์ เทียบกับ มีนาคม 2026)"
        period_thai_para = period_thai_frame.paragraphs[0]
        period_thai_para.font.size = Pt(16)
        period_thai_para.font.color.rgb = self.COLORS['medium_gray']
        period_thai_para.alignment = PP_ALIGN.CENTER

        # Bottom decoration
        bottom_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(5.325),
            self.prs.slide_width,
            Inches(0.3)
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = self.COLORS['secondary_blue']
        bottom_bar.line.fill.background()

    def add_executive_summary(self):
        """Executive summary with key highlights"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['white']
        bg.line.fill.background()

        # Header
        self._add_slide_header(slide, "Executive Summary", "(สรุปผู้บริหาร)")

        # Key metrics boxes
        metrics = [
            {
                'title': 'February Achievement\n(ผลงานเดือนกุมภาพันธ์)',
                'value': '73.86%',
                'subtitle': 'of Target\n(ของเป้าหมาย)',
                'color': self.COLORS['accent_green'],
                'x': 0.5
            },
            {
                'title': 'March Achievement\n(ผลงานเดือนมีนาคม)',
                'value': '56.97%',
                'subtitle': 'of Target\n(ของเป้าหมาย)',
                'color': self.COLORS['accent_red'],
                'x': 3.5
            },
            {
                'title': 'Average Transaction\n(ค่าเฉลี่ยต่อบิล)',
                'value': '+1.4%',
                'subtitle': 'Month-over-Month\n(เทียบเดือนก่อน)',
                'color': self.COLORS['primary_blue'],
                'x': 6.5
            }
        ]

        for metric in metrics:
            self._add_metric_box(slide, metric)

        # Key observations
        obs_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5),
                                           Inches(9), Inches(1.5))
        obs_frame = obs_box.text_frame
        obs_frame.word_wrap = True

        observations = [
            "• February benefited from Chinese New Year & Valentine's Day festivities",
            "  (กุมภาพันธ์ได้รับประโยชน์จากเทศกาลตรุษจีนและวาเลนไทน์)",
            "",
            "• March promotions showed lower customer interest due to repetitive giveaways",
            "  (โปรโมชั่นเดือนมีนาคมมีความสนใจน้อยลงเนื่องจากของแถมซ้ำซาก)",
            "",
            "• Economic factors and pre-Songkran savings affected spending patterns",
            "  (ปัจจัยทางเศรษฐกิจและการออมเงินก่อนสงกรานต์ส่งผลต่อรูปแบบการใช้จ่าย)"
        ]

        for i, obs in enumerate(observations):
            p = obs_frame.add_paragraph() if i > 0 else obs_frame.paragraphs[0]
            p.text = obs
            p.font.size = Pt(11)
            p.font.color.rgb = self.COLORS['dark_gray']
            p.space_after = Pt(3)

    def add_february_performance(self):
        """February 2026 detailed performance"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['white']
        bg.line.fill.background()

        # Header
        self._add_slide_header(slide, "February 2026 Performance",
                              "(ผลงานเดือนกุมภาพันธ์ 2026)")

        # Sales metrics
        metrics_data = [
            ("Target", "เป้าหมาย", "฿1,070,000", Inches(0.5)),
            ("Actual Sales", "ยอดขายจริง", "฿790,319.50", Inches(3)),
            ("Achievement", "ทำได้", "73.86%", Inches(5.5)),
            ("Store Traffic", "ลูกค้าเข้าร้าน", "1,304", Inches(8))
        ]

        for title_en, title_th, value, x_pos in metrics_data:
            self._add_data_card(slide, x_pos, Inches(1.2), title_en, title_th, value)

        # Transaction metrics
        trans_data = [
            ("Total Items Sold", "จำนวนสินค้าที่ขายได้", "2,675 pcs", Inches(0.5)),
            ("Total Bills", "จำนวนบิล", "552 bills", Inches(2.75)),
            ("ATV", "ค่าเฉลี่ยต่อบิล", "฿1,431.74", Inches(5)),
            ("UPT", "ชิ้นต่อบิล", "4.85 pcs", Inches(7.25))
        ]

        for title_en, title_th, value, x_pos in trans_data:
            self._add_data_card(slide, x_pos, Inches(2.8), title_en, title_th, value)

        # Promotions section
        promo_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2),
                                             Inches(9), Inches(1.2))
        promo_frame = promo_box.text_frame
        promo_frame.word_wrap = True

        # Promo title
        promo_title = promo_frame.paragraphs[0]
        promo_title.text = "Key Promotions (โปรโมชั่นหลัก)"
        promo_title.font.size = Pt(14)
        promo_title.font.bold = True
        promo_title.font.color.rgb = self.COLORS['primary_blue']

        promos = [
            "✓ Free Pencil Pouch (฿1,500 purchase) | ✓ Free Pods (฿2,000 purchase)",
            "✓ Special Doodle Engraving + Name (Bundle set + Boot purchase)",
            "→ High interest due to cute Pods design and exclusive engraving feature"
        ]

        for promo in promos:
            p = promo_frame.add_paragraph()
            p.text = promo
            p.font.size = Pt(10)
            p.font.color.rgb = self.COLORS['dark_gray']
            p.space_after = Pt(2)

    def add_march_performance(self):
        """March 2026 detailed performance"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['white']
        bg.line.fill.background()

        # Header
        self._add_slide_header(slide, "March 2026 Performance (1-29)",
                              "(ผลงานเดือนมีนาคม 2026)")

        # Sales metrics
        metrics_data = [
            ("Actual Sales", "ยอดขายจริง", "฿548,598.75", Inches(0.5)),
            ("Achievement", "ทำได้", "56.97%", Inches(3)),
            ("Store Traffic", "ลูกค้าเข้าร้าน", "1,074", Inches(5.5))
        ]

        for title_en, title_th, value, x_pos in metrics_data:
            self._add_data_card(slide, x_pos, Inches(1.2), title_en, title_th, value,
                              color=self.COLORS['accent_red'])

        # Transaction metrics
        trans_data = [
            ("Total Items Sold", "จำนวนสินค้าที่ขายได้", "1,889 pcs", Inches(0.5)),
            ("Total Bills", "จำนวนบิล", "378 bills", Inches(2.75)),
            ("ATV", "ค่าเฉลี่ยต่อบิล", "฿1,451.32", Inches(5)),
            ("UPT", "ชิ้นต่อบิล", "5.0 pcs", Inches(7.25))
        ]

        for title_en, title_th, value, x_pos in trans_data:
            self._add_data_card(slide, x_pos, Inches(2.8), title_en, title_th, value)

        # Promotions section
        promo_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2),
                                             Inches(9), Inches(1.2))
        promo_frame = promo_box.text_frame
        promo_frame.word_wrap = True

        # Promo title
        promo_title = promo_frame.paragraphs[0]
        promo_title.text = "Key Promotions (โปรโมชั่นหลัก)"
        promo_title.font.size = Pt(14)
        promo_title.font.bold = True
        promo_title.font.color.rgb = self.COLORS['primary_blue']

        promos = [
            "Free Jelly Pouch/Shop Tote/Clean Lens (฿1,800 purchase)",
            "Free Boot+ (32oz Bundle set purchase, Mar 20 - Apr 19)",
            "⚠ Lower interest: repetitive giveaways & higher threshold vs. average spend (฿1,400)"
        ]

        for promo in promos:
            p = promo_frame.add_paragraph()
            p.text = promo
            p.font.size = Pt(10)
            p.font.color.rgb = self.COLORS['dark_gray']
            p.space_after = Pt(2)

    def add_comparison_chart(self):
        """Visual comparison chart"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['white']
        bg.line.fill.background()

        # Header
        self._add_slide_header(slide, "Month-over-Month Comparison",
                              "(เปรียบเทียบรายเดือน)")

        # Create clustered column chart
        chart_data = CategoryChartData()
        chart_data.categories = ['Sales\n(ยอดขาย)', 'Items\n(สินค้า)', 'Bills\n(บิล)',
                                 'Traffic\n(ลูกค้า)', 'ATV\n(ค่าเฉลี่ย)']

        # Normalize data for visualization (using percentage of max)
        feb_data = [790319.50, 2675, 552, 1304, 1431.74]
        mar_data = [548598.75, 1889, 378, 1074, 1451.32]

        # Calculate percentages for better visualization
        max_values = [max(feb_data[i], mar_data[i]) for i in range(5)]
        feb_pct = [(feb_data[i] / max_values[i] * 100) for i in range(5)]
        mar_pct = [(mar_data[i] / max_values[i] * 100) for i in range(5)]

        chart_data.add_series('February (กุมภาพันธ์)', tuple(feb_pct))
        chart_data.add_series('March (มีนาคม)', tuple(mar_pct))

        # Add chart
        x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(3.2)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        # Style chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)

        # Configure value axis (y-axis)
        value_axis = chart.value_axis
        value_axis.minimum_scale = 0
        value_axis.maximum_scale = 100
        value_axis.major_unit = 20
        value_axis.has_major_gridlines = True
        value_axis.tick_labels.font.size = Pt(9)

        # Add y-axis title
        value_axis.has_title = True
        value_axis.axis_title.text_frame.text = "Index (100 = Highest Value)\n(ดัชนี: 100 = ค่าสูงสุด)"
        value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
        value_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = self.COLORS['dark_gray']

        # Configure category axis (x-axis)
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(9)

        # Format chart colors
        series_feb = chart.series[0]
        series_feb.format.fill.solid()
        series_feb.format.fill.fore_color.rgb = self.COLORS['primary_blue']

        series_mar = chart.series[1]
        series_mar.format.fill.solid()
        series_mar.format.fill.fore_color.rgb = self.COLORS['secondary_blue']

        # Add explanation note
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(5),
                                           Inches(9), Inches(0.5))
        note_frame = note_box.text_frame
        note_frame.word_wrap = True

        note_text = note_frame.paragraphs[0]
        note_text.text = "Note: Values are normalized to 100 (highest value in each category) for visual comparison across different metrics. " \
                        "Each metric uses its own scale. (หมายเหตุ: ค่าต่างๆ ถูกปรับให้เป็น 100 สำหรับการเปรียบเทียบ แต่ละตัวชี้วัดใช้มาตราส่วนของตัวเอง)"
        note_text.font.size = Pt(8)
        note_text.font.italic = True
        note_text.font.color.rgb = self.COLORS['medium_gray']

    def add_insights_recommendations(self):
        """Key insights and recommendations"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['white']
        bg.line.fill.background()

        # Header
        self._add_slide_header(slide, "Key Insights & Recommendations",
                              "(ข้อมูลเชิงลึกและข้อเสนอแนะ)")

        # Left column - Insights
        insights_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3),
                                               Inches(4.5), Inches(3.8))
        insights_frame = insights_box.text_frame
        insights_frame.word_wrap = True

        insights_title = insights_frame.paragraphs[0]
        insights_title.text = "Key Insights (ข้อมูลเชิงลึก)"
        insights_title.font.size = Pt(16)
        insights_title.font.bold = True
        insights_title.font.color.rgb = self.COLORS['primary_blue']
        insights_title.space_after = Pt(10)

        insights = [
            ("🎯 Festival Impact", "ผลกระทบจากเทศกาล",
             "Chinese New Year & Valentine's drove 17% higher achievement in Feb"),

            ("📊 Promotion Fatigue", "ความเบื่อหน่ายจากโปรโมชั่น",
             "Repetitive giveaways (shop tote, jelly pouch) reduced March interest"),

            ("💰 Price Threshold", "เกณฑ์ราคา",
             "฿1,800 threshold exceeds avg. spend (฿1,400), limiting conversion"),

            ("🌍 External Factors", "ปัจจัยภายนอก",
             "Economic conditions, war concerns, fuel prices affected spending"),

            ("✈️ Pre-Songkran Saving", "การออมก่อนสงกรานต์",
             "Customers saving for holiday travel reduced discretionary spending")
        ]

        for emoji, title_en, desc in insights:
            p = insights_frame.add_paragraph()
            p.text = f"{emoji} {title_en}"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['dark_gray']
            p.space_after = Pt(2)

            p2 = insights_frame.add_paragraph()
            p2.text = f"   ({desc})"
            p2.font.size = Pt(9)
            p2.font.color.rgb = self.COLORS['medium_gray']
            p2.space_after = Pt(8)

        # Right column - Recommendations
        rec_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3),
                                           Inches(4.3), Inches(3.8))
        rec_frame = rec_box.text_frame
        rec_frame.word_wrap = True

        rec_title = rec_frame.paragraphs[0]
        rec_title.text = "Recommendations (ข้อเสนอแนะ)"
        rec_title.font.size = Pt(16)
        rec_title.font.bold = True
        rec_title.font.color.rgb = self.COLORS['accent_green']
        rec_title.space_after = Pt(10)

        recommendations = [
            ("🎁 Fresh Incentives", "สิ่งจูงใจใหม่",
             "Introduce new, desirable giveaways to reignite interest"),

            ("💵 Tiered Pricing", "ราคาแบบขั้นบันได",
             "Create ฿1,200-1,500-1,800 tiers matching customer spend patterns"),

            ("📦 Product Mix", "การผสมผสานสินค้า",
             "Add merch/accessories for easy upselling to threshold"),

            ("🎪 Event Marketing", "การตลาดเชิงกิจกรรม",
             "Create micro-events to drive foot traffic during slow periods"),

            ("👥 Target 16-25oz", "กลุ่มเป้าหมาย 16-25 ออนซ์",
             "Focus promotions on popular sizes for office/female demographic"),

            ("📱 Social Campaigns", "แคมเปญโซเชียล",
             "Leverage exclusive engraving/designs that performed well in Feb")
        ]

        for emoji, title_en, desc in recommendations:
            p = rec_frame.add_paragraph()
            p.text = f"{emoji} {title_en}"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['dark_gray']
            p.space_after = Pt(2)

            p2 = rec_frame.add_paragraph()
            p2.text = f"   ({desc})"
            p2.font.size = Pt(9)
            p2.font.color.rgb = self.COLORS['medium_gray']
            p2.space_after = Pt(8)

    def add_action_plan(self):
        """Action plan and next steps"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background with gradient effect
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['light_gray']
        bg.line.fill.background()

        # Header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                              Inches(9), Inches(0.8))
        header_frame = header_box.text_frame
        header_frame.text = "Action Plan for April 2026"
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(32)
        header_para.font.bold = True
        header_para.font.color.rgb = self.COLORS['primary_blue']
        header_para.alignment = PP_ALIGN.CENTER

        # Thai subtitle
        thai_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8),
                                            Inches(9), Inches(0.4))
        thai_frame = thai_box.text_frame
        thai_frame.text = "(แผนปฏิบัติการสำหรับเดือนเมษายน 2026)"
        thai_para = thai_frame.paragraphs[0]
        thai_para.font.size = Pt(16)
        thai_para.font.color.rgb = self.COLORS['medium_gray']
        thai_para.alignment = PP_ALIGN.CENTER

        # Action items with timeline
        actions = [
            {
                'phase': 'Week 1-2 (สัปดาห์ 1-2)',
                'title': 'Promotion Design',
                'items': [
                    '• Design new exclusive giveaway items based on Feb success',
                    '• Create tiered promotion structure (฿1,200/฿1,500/฿1,800)',
                    '• Source new merch products for upselling opportunities'
                ],
                'color': self.COLORS['primary_blue']
            },
            {
                'phase': 'Week 2-3 (สัปดาห์ 2-3)',
                'title': 'Marketing Campaign',
                'items': [
                    '• Launch social media teaser campaign for new promotions',
                    '• Plan Songkran-themed special events and activities',
                    '• Prepare staff training on new upselling techniques'
                ],
                'color': self.COLORS['secondary_blue']
            },
            {
                'phase': 'Week 3-4 (สัปดาห์ 3-4)',
                'title': 'Execution & Monitoring',
                'items': [
                    '• Launch new promotion campaign',
                    '• Monitor daily sales vs. target with adjusted benchmarks',
                    '• Gather customer feedback and adjust tactics in real-time'
                ],
                'color': self.COLORS['accent_green']
            }
        ]

        y_position = 1.5
        for action in actions:
            # Phase box
            phase_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.7), Inches(y_position),
                Inches(8.6), Inches(1)
            )
            phase_box.fill.solid()
            phase_box.fill.fore_color.rgb = action['color']
            phase_box.line.color.rgb = action['color']
            phase_box.shadow.inherit = False

            # Phase title
            phase_text = phase_box.text_frame
            phase_title = phase_text.paragraphs[0]
            phase_title.text = f"{action['phase']}: {action['title']}"
            phase_title.font.size = Pt(13)
            phase_title.font.bold = True
            phase_title.font.color.rgb = self.COLORS['white']
            phase_title.alignment = PP_ALIGN.LEFT
            phase_text.margin_left = Inches(0.2)
            phase_text.margin_top = Inches(0.1)

            # Action items
            for item in action['items']:
                p = phase_text.add_paragraph()
                p.text = item
                p.font.size = Pt(10)
                p.font.color.rgb = self.COLORS['white']
                p.level = 0
                p.space_after = Pt(2)

            y_position += 1.2

    def _add_slide_header(self, slide, title_en, title_th):
        """Add consistent slide header"""
        # Top accent bar
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.prs.slide_width,
            Inches(0.15)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.COLORS['primary_blue']
        accent.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25),
                                             Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title_en
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['primary_blue']

        # Thai subtitle
        thai_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.65),
                                            Inches(9), Inches(0.3))
        thai_frame = thai_box.text_frame
        thai_frame.text = title_th
        thai_para = thai_frame.paragraphs[0]
        thai_para.font.size = Pt(12)
        thai_para.font.color.rgb = self.COLORS['medium_gray']

    def _add_metric_box(self, slide, metric):
        """Add metric display box"""
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(metric['x']), Inches(1.5),
            Inches(2.5), Inches(1.3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = self.COLORS['white']
        box.line.color.rgb = metric['color']
        box.line.width = Pt(2)

        text_frame = box.text_frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Title
        title = text_frame.paragraphs[0]
        title.text = metric['title']
        title.font.size = Pt(10)
        title.font.color.rgb = self.COLORS['dark_gray']
        title.alignment = PP_ALIGN.CENTER
        title.space_after = Pt(5)

        # Value
        value_p = text_frame.add_paragraph()
        value_p.text = metric['value']
        value_p.font.size = Pt(20)
        value_p.font.bold = True
        value_p.font.color.rgb = metric['color']
        value_p.alignment = PP_ALIGN.CENTER
        value_p.space_after = Pt(3)

        # Subtitle
        sub_p = text_frame.add_paragraph()
        sub_p.text = metric['subtitle']
        sub_p.font.size = Pt(8)
        sub_p.font.color.rgb = self.COLORS['medium_gray']
        sub_p.alignment = PP_ALIGN.CENTER

    def _add_data_card(self, slide, x, y, title_en, title_th, value,
                       color=None):
        """Add data display card"""
        if color is None:
            color = self.COLORS['primary_blue']

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y,
            Inches(2), Inches(1.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.COLORS['light_gray']
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        text_frame = card.text_frame
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # English title
        title_p = text_frame.paragraphs[0]
        title_p.text = title_en
        title_p.font.size = Pt(10)
        title_p.font.bold = True
        title_p.font.color.rgb = self.COLORS['dark_gray']
        title_p.alignment = PP_ALIGN.CENTER
        title_p.space_after = Pt(2)

        # Thai title
        thai_p = text_frame.add_paragraph()
        thai_p.text = f"({title_th})"
        thai_p.font.size = Pt(8)
        thai_p.font.color.rgb = self.COLORS['medium_gray']
        thai_p.alignment = PP_ALIGN.CENTER
        thai_p.space_after = Pt(8)

        # Value
        value_p = text_frame.add_paragraph()
        value_p.text = value
        value_p.font.size = Pt(16)
        value_p.font.bold = True
        value_p.font.color.rgb = color
        value_p.alignment = PP_ALIGN.CENTER

    def save(self, filename="SUNA_Sales_Analysis_Feb_Mar_2026.pptx"):
        """Save presentation to file"""
        filepath = os.path.join("/teamspace/studios/this_studio/comprehensive-suna-bim-agent",
                               filename)
        self.prs.save(filepath)
        return filepath

def main():
    """Generate complete sales presentation"""
    print("🎨 Creating SUNA Sales Analysis Presentation...")

    presentation = SalesPresentation()

    print("  → Adding title slide...")
    presentation.add_title_slide()

    print("  → Adding executive summary...")
    presentation.add_executive_summary()

    print("  → Adding February 2026 performance...")
    presentation.add_february_performance()

    print("  → Adding March 2026 performance...")
    presentation.add_march_performance()

    print("  → Adding comparison chart...")
    presentation.add_comparison_chart()

    print("  → Adding insights & recommendations...")
    presentation.add_insights_recommendations()

    print("  → Adding action plan...")
    presentation.add_action_plan()

    filepath = presentation.save()
    print(f"\n✅ Presentation created successfully!")
    print(f"📁 Location: {filepath}")
    print(f"📊 Total slides: 7")
    print(f"🎯 Ready for executive review")

    return filepath

if __name__ == "__main__":
    main()
